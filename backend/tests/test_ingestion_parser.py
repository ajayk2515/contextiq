from pathlib import Path

import pytest
import tiktoken
from docx import Document as WordDocument
from pptx import Presentation
from reportlab.pdfgen.canvas import Canvas

from app.ingestion.parser import parse_and_chunk


def create_sample(path: Path) -> None:
    if path.suffix == ".md":
        path.write_text("# Leave Policy\n\nEmployees receive twenty days of annual leave.")
    elif path.suffix == ".docx":
        document = WordDocument()
        document.add_heading("Leave Policy", level=1)
        document.add_paragraph("Employees receive twenty days of annual leave.")
        document.save(path)
    elif path.suffix == ".pptx":
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Leave Policy"
        slide.placeholders[1].text = "Employees receive twenty days of annual leave."
        presentation.save(path)
    elif path.suffix == ".pdf":
        canvas = Canvas(str(path))
        canvas.drawString(72, 760, "Leave Policy")
        canvas.drawString(72, 730, "Employees receive twenty days of annual leave.")
        canvas.save()


@pytest.mark.parametrize("extension", [".pdf", ".docx", ".pptx", ".md"])
def test_docling_parses_required_formats(tmp_path: Path, extension: str) -> None:
    path = tmp_path / f"policy{extension}"
    create_sample(path)

    chunks = parse_and_chunk(path, chunk_size=800, overlap=120)

    assert chunks
    assert any("annual leave" in chunk.text.lower() for chunk in chunks)
    assert all(len(chunk.chunk_hash) == 64 for chunk in chunks)


def test_chunking_preserves_headers_and_configured_overlap(tmp_path: Path) -> None:
    first = " ".join(f"alpha{i}" for i in range(350))
    second = " ".join(f"beta{i}" for i in range(350))
    path = tmp_path / "long-policy.md"
    path.write_text(f"# Policy\n\n{first}\n\n## Exceptions\n\n{second}")

    chunks = parse_and_chunk(path, chunk_size=200, overlap=30)
    encoding = tiktoken.get_encoding("cl100k_base")

    assert len(chunks) > 2
    assert all(len(encoding.encode(chunk.text)) <= 200 for chunk in chunks)
    assert any(chunk.section for chunk in chunks)
    assert len({chunk.chunk_hash for chunk in chunks}) == len(chunks)
